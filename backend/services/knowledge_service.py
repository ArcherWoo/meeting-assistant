"""
知识库服务 - 文件导入 Pipeline + LanceDB 向量管理
遵循 PRD §12 知识库与 RAG 引擎：
  文件上传 → 解析 → LLM 结构化提取 → SQLite 写入 → 文本分块 → Embedding → LanceDB 写入
支持: PPT/PPTX, PDF, DOC/DOCX, 图片, 纯文本
"""
import json
import io
import logging
import os
import shutil
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Optional, Any

from services.document_parsing import chunk_parsed_document, document_parser_registry
from services.document_parsing.models import ParsedDocument
from services.document_parsing.prompt_render import render_document_for_prompt
from services.runtime_controls import attachment_parse_controller
from services.runtime_paths import IMPORTED_FILES_DIR, VECTORS_DIR
from services.storage import storage, gen_id
from utils.decryption_handler import decrypt_esafenet_file

logger = logging.getLogger(__name__)

# LLM 结构化提取 Prompt（PRD §12.4）
EXTRACTION_PROMPT = """请从以下 PPT 内容中提取所有采购相关的结构化信息。
对每个采购项，提取以下字段并以 JSON 数组格式返回：
[
  {
    "category": "采购品类（如办公设备、IT设备等）",
    "item_name": "具体品名",
    "supplier": "供应商名称",
    "unit_price": 数字或null,
    "quantity": 数字或null,
    "total_price": 数字或null,
    "currency": "CNY",
    "delivery_days": 数字或null,
    "payment_terms": "付款方式",
    "procurement_date": "日期或null",
    "raw_text": "原始文本片段（用于溯源）"
  }
]

如果没有找到采购相关信息，返回空数组 []。
只返回 JSON，不要其他文字。

PPT 内容：
{content}"""


class KnowledgeService:
    """知识库服务 - 管理文件导入和向量检索"""

    # PPT 扩展名
    PPT_EXTS = {".ppt", ".pptx"}
    # 纯文本扩展名
    TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".xml"}
    # 图片扩展名
    IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}

    FAST_EXCEL_MAX_ROWS_PER_SHEET = 200
    FAST_EXCEL_MAX_COLS = 20
    FAST_PPT_MAX_TABLE_ROWS = 20

    def __init__(self) -> None:
        self._lance_db = None
        self._chunks_table = None

    @staticmethod
    def _build_imported_file_path(import_id: str, filename: str) -> Path:
        safe_name = Path(filename or "uploaded-file").name or "uploaded-file"
        return (IMPORTED_FILES_DIR / import_id / safe_name).resolve()

    async def _persist_imported_file(self, import_id: str, filename: str, file_content: bytes) -> Path:
        target_path = self._build_imported_file_path(import_id, filename)

        def _write_file() -> Path:
            target_path.parent.mkdir(parents=True, exist_ok=True)
            target_path.write_bytes(file_content)
            return target_path

        written_path = await asyncio.to_thread(_write_file)
        await storage.update_ppt_import_file_path(import_id, str(written_path))
        return written_path

    @staticmethod
    def _build_chunk_metadata(import_id: str, chunk: dict) -> str:
        """将片段定位信息编码到已有 metadata 字段，避免变更 LanceDB 主表结构。"""
        metadata = {"import_id": import_id}
        raw_metadata = chunk.get("metadata_json")
        if isinstance(raw_metadata, str) and raw_metadata.strip():
            try:
                parsed = json.loads(raw_metadata)
                if isinstance(parsed, dict):
                    metadata.update(parsed)
            except json.JSONDecodeError:
                pass
        for key in ("chunk_index", "char_start", "char_end"):
            value = chunk.get(key)
            if value is not None:
                metadata[key] = value
        return json.dumps(metadata, ensure_ascii=False)

    @staticmethod
    def _parse_chunk_metadata(raw_metadata: Any) -> dict:
        """兼容历史字符串 metadata 和未来可能的 dict 形态。"""
        if isinstance(raw_metadata, dict):
            return raw_metadata
        if not raw_metadata:
            return {}
        if isinstance(raw_metadata, str):
            try:
                parsed = json.loads(raw_metadata)
                return parsed if isinstance(parsed, dict) else {}
            except json.JSONDecodeError:
                return {}
        return {}

    @staticmethod
    def _escape_lancedb_literal(value: str) -> str:
        return value.replace("'", "''")

    @staticmethod
    def _extract_pdf_text_with_pymupdf(file_content: bytes) -> str:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=file_content, filetype="pdf")
        try:
            parts = []
            for page in doc:
                text = page.get_text().strip()
                if text:
                    parts.append(text)
            return "\n\n".join(parts)
        finally:
            doc.close()

    @staticmethod
    def _extract_pdf_text_with_pypdf(file_content: bytes) -> str:
        try:
            from pypdf import PdfReader
        except ModuleNotFoundError:
            from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(file_content))
        parts = []
        for page in reader.pages:
            text = (page.extract_text() or "").strip()
            if text:
                parts.append(text)
        return "\n\n".join(parts)

    def _extract_pdf_text_sync(self, file_content: bytes, filename: str) -> str:
        parsers = (
            ("PyMuPDF", self._extract_pdf_text_with_pymupdf),
            ("pypdf", self._extract_pdf_text_with_pypdf),
        )
        installed_parser_names: list[str] = []
        parser_errors: list[str] = []

        for parser_name, parser in parsers:
            try:
                text = parser(file_content).strip()
                installed_parser_names.append(parser_name)
                if text:
                    return text
                logger.warning(f"{parser_name} 未能从 PDF 提取到文本: {filename}")
            except ModuleNotFoundError:
                continue
            except ImportError:
                continue
            except Exception as e:
                installed_parser_names.append(parser_name)
                logger.warning(f"{parser_name} 解析 PDF 失败 {filename}: {e}")
                parser_errors.append(f"{parser_name}: {e}")

        if installed_parser_names:
            if parser_errors:
                details = "; ".join(parser_errors)
                return f"[PDF 文件: {filename}，解析失败: {details}]"
            return f"[PDF 文件: {filename}，未提取到可读文本，可能是扫描件或图片型 PDF]"

        logger.warning(f"PDF 解析依赖缺失: {filename}")
        return f"[PDF 文件: {filename}，缺少 PDF 解析依赖，请安装 pypdf 或 PyMuPDF]"

    def _extract_ppt_text_fast_sync(self, file_content: bytes, filename: str) -> str:
        try:
            from pptx import Presentation
        except Exception as e:
            return f"[PPT 文件: {filename}，缺少 python-pptx 依赖: {e}]"

        try:
            presentation = Presentation(io.BytesIO(file_content))
        except Exception as e:
            return f"[PPT 文件: {filename}，解析失败: {e}]"

        sections: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []
            title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else ""
            slide_lines.append(title or f"幻灯片 {slide_index}")

            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text.strip()
                    if text and text != title:
                        slide_lines.append(text)
                elif getattr(shape, "has_table", False):
                    for row_index, row in enumerate(shape.table.rows, start=1):
                        if row_index > self.FAST_PPT_MAX_TABLE_ROWS:
                            slide_lines.append("[表格内容已截断，仅保留前 20 行]")
                            break
                        compact_cells = [
                            cell.text.strip()
                            for cell in row.cells[: self.FAST_EXCEL_MAX_COLS]
                            if cell.text and cell.text.strip()
                        ]
                        if compact_cells:
                            slide_lines.append(" | ".join(compact_cells))

            compact_lines = [line for line in slide_lines if line]
            if compact_lines:
                sections.append("\n".join([f"## {compact_lines[0]}", *compact_lines[1:]]))

        return "\n\n".join(sections)

    def _extract_excel_text_fast_sync(self, file_content: bytes, filename: str) -> str:
        try:
            import openpyxl
        except Exception as e:
            return f"[Excel 文件: {filename}，缺少 openpyxl 依赖: {e}]"

        try:
            workbook = openpyxl.load_workbook(
                io.BytesIO(file_content),
                read_only=True,
                data_only=True,
            )
        except Exception as e:
            return f"[Excel 文件: {filename}，解析失败: {e}]"

        try:
            sections: list[str] = []
            for sheet in workbook.worksheets:
                lines: list[str] = [f"## Sheet: {sheet.title}"]
                non_empty_rows = 0
                for row in sheet.iter_rows(max_col=self.FAST_EXCEL_MAX_COLS, values_only=True):
                    cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                    if not cells:
                        continue
                    non_empty_rows += 1
                    if non_empty_rows > self.FAST_EXCEL_MAX_ROWS_PER_SHEET:
                        lines.append("[表格内容已截断，仅保留前 200 行非空数据]")
                        break
                    lines.append(" | ".join(cells))
                sections.append("\n".join(lines))
            return "\n\n".join(sections)
        finally:
            workbook.close()

    async def _delete_vectors_for_file(self, filename: str) -> bool:
        if not self._lance_db or "doc_chunks" not in (self._lance_db.table_names() or []):
            return False
        try:
            table = self._lance_db.open_table("doc_chunks")
            safe_filename = self._escape_lancedb_literal(filename)
            table.delete(f"source_file = '{safe_filename}'")
            return True
        except Exception as e:
            logger.warning(f"删除向量失败: {e}")
            return False

    async def _persist_chunks(
        self, chunks: List[dict], import_id: str, filename: str, file_type: str,
    ) -> int:
        if not chunks:
            return 0
        return await storage.add_knowledge_chunks(
            import_id=import_id,
            source_file=filename,
            file_type=file_type,
            chunks=chunks,
        )

    async def _reset_import_index(
        self, import_id: str, filename: str, clear_procurement: bool = False,
    ) -> None:
        await storage.delete_knowledge_chunks(import_id=import_id)
        if clear_procurement:
            await storage.db.execute(
                "DELETE FROM procurement_records WHERE source_file_id=?",
                (import_id,),
            )
            await storage.db.commit()
        await self._delete_vectors_for_file(filename)

    async def initialize(self) -> None:
        """初始化 LanceDB 连接（异步，避免阻塞事件循环）"""
        import asyncio
        VECTORS_DIR.mkdir(parents=True, exist_ok=True)
        try:
            import lancedb
            # lancedb.connect 是同步阻塞调用，必须放入线程池以防阻塞事件循环
            self._lance_db = await asyncio.to_thread(lancedb.connect, str(VECTORS_DIR))
            logger.info(f"LanceDB 已连接: {VECTORS_DIR}")
        except ImportError:
            logger.warning("lancedb 未安装，向量检索功能不可用")
        except Exception as e:
            logger.warning(f"LanceDB 初始化失败: {e}")

    async def initialize(self) -> None:
        """初始化 LanceDB 连接。启动阶段超时后降级为无向量连接模式。"""
        import asyncio

        VECTORS_DIR.mkdir(parents=True, exist_ok=True)
        timeout_raw = os.getenv("MEETING_ASSISTANT_LANCEDB_INIT_TIMEOUT_SEC", "").strip()
        try:
            init_timeout_sec = max(1.0, float(timeout_raw)) if timeout_raw else 5.0
        except ValueError:
            init_timeout_sec = 5.0

        try:
            import lancedb

            self._lance_db = await asyncio.wait_for(
                asyncio.to_thread(lancedb.connect, str(VECTORS_DIR)),
                timeout=init_timeout_sec,
            )
            logger.info(f"LanceDB 已连接: {VECTORS_DIR}")
        except asyncio.TimeoutError:
            self._lance_db = None
            logger.warning(
                "LanceDB 初始化超时（%ss），本次启动先跳过向量库连接，服务继续启动",
                init_timeout_sec,
            )
        except ImportError:
            logger.warning("lancedb 未安装，向量检索功能不可用")
        except Exception as e:
            logger.warning(f"LanceDB 初始化失败: {e}")

    async def ingest_file(
        self, file_content: bytes, filename: str,
        llm_fn: Optional[object] = None,
        embedding_fn: Optional[object] = None,
        owner_id: Optional[str] = None,
        is_encrypted: bool = False,
    ) -> dict:
        """
        通用文件导入入口 - 根据文件类型分发到对应的处理器
        若 is_encrypted=True，先通过解密 hook 处理文件字节。
        """
        if is_encrypted:
            file_content = decrypt_esafenet_file(file_content, filename)

        parsed_document = await self._parse_document(file_content, filename)
        return await self._ingest_parsed_document(
            file_content=file_content,
            filename=filename,
            parsed_document=parsed_document,
            llm_fn=llm_fn,
            embedding_fn=embedding_fn,
            owner_id=owner_id,
        )

    async def ingest_ppt(
        self, file_content: bytes, filename: str,
        llm_fn: Optional[object] = None,
        embedding_fn: Optional[object] = None,
        owner_id: Optional[str] = None,
    ) -> dict:
        """
        PPT 导入知识库完整 Pipeline
        返回: {"import_id": str, "status": str, "extracted_count": int, "chunks_count": int}
        """
        from services.ppt_parser import PPTParser
        import hashlib

        # 1. 文件去重检查
        file_hash = hashlib.md5(file_content).hexdigest()
        import_id = await storage.record_ppt_import(
            file_name=filename, file_hash=file_hash,
            file_size=len(file_content), slide_count=0,
            owner_id=owner_id,
        )
        # 检查是否已导入（record_ppt_import 返回已有 ID）
        existing = await storage._fetchone(
            "SELECT import_status, slide_count FROM ppt_imports WHERE id=?", (import_id,)
        )
        existing_chunk_count = await storage.count_knowledge_chunks(import_id=import_id)
        if existing and existing["import_status"] == "completed" and existing_chunk_count > 0:
            # 即使是重复导入，也重新解析以获取统计信息
            try:
                parser = PPTParser()
                ppt_data = await parser.parse(file_content, filename)
                dup_slide_count = len(ppt_data.get("slides", []))
                dup_text_length = len(ppt_data.get("full_markdown", ""))
                dup_table_count = ppt_data.get("extraction_stats", {}).get("total_tables", 0)
                dup_image_count = ppt_data.get("extraction_stats", {}).get("total_images", 0)
                dup_chunks = len(self._split_into_chunks(ppt_data, filename))
            except Exception:
                dup_slide_count = existing.get("slide_count", 0) or 0
                dup_text_length = 0
                dup_table_count = 0
                dup_image_count = 0
                dup_chunks = 0
            return {
                "import_id": import_id, "status": "duplicate",
                "file_type": "ppt",
                "slide_count": dup_slide_count, "text_length": dup_text_length,
                "table_count": dup_table_count, "image_count": dup_image_count,
                "extracted_count": 0, "chunks_count": existing_chunk_count or dup_chunks,
            }

        await storage.update_ppt_import_status(import_id, "processing")

        # 2. 解析 PPT
        parser = PPTParser()
        ppt_data = await parser.parse(file_content, filename)
        slide_count = len(ppt_data.get("slides", []))
        full_markdown = ppt_data.get("full_markdown", "")
        text_length = len(full_markdown)
        table_count = ppt_data.get("extraction_stats", {}).get("total_tables", 0)
        image_count = ppt_data.get("extraction_stats", {}).get("total_images", 0)
        text_chunks = self._split_into_chunks(ppt_data, filename)
        chunks_parsed = len(text_chunks)

        # 旧版本重复导入时可能只有导入记录，没有可检索正文；这里统一重建索引。
        await self._reset_import_index(import_id, filename, clear_procurement=True)
        await storage.db.execute(
            "UPDATE ppt_imports SET slide_count=? WHERE id=?", (slide_count, import_id)
        )
        await storage.db.commit()

        # 3. LLM 结构化提取（如果提供了 LLM 函数）
        extracted_count = 0
        if llm_fn and full_markdown:
            extracted_count = await self._extract_procurement_fields(
                full_markdown, filename, import_id, llm_fn
            )

        # 4. 文本分块 + Embedding + 写入 LanceDB
        stored_chunks_count = await self._persist_chunks(
            text_chunks, import_id, filename, "ppt",
        )
        vector_chunks_count = 0
        if embedding_fn and self._lance_db:
            vector_chunks_count = await self._vectorize_chunks(
                text_chunks, import_id, embedding_fn
            )

        await storage.update_ppt_import_status(import_id, "completed", extracted_count)

        logger.info(
            f"PPT 导入完成: {filename}, {slide_count} 页, {text_length} 字符, "
            f"{stored_chunks_count} 个正文块, {vector_chunks_count} 个向量块"
        )

        return {
            "import_id": import_id,
            "status": "completed",
            "file_type": "ppt",
            "slide_count": slide_count,
            "text_length": text_length,
            "table_count": table_count,
            "image_count": image_count,
            "extracted_count": extracted_count,
            "stored_chunks_count": stored_chunks_count,
            "vector_chunks_count": vector_chunks_count,
            "chunks_count": vector_chunks_count if vector_chunks_count else stored_chunks_count or chunks_parsed,
        }

    async def _extract_procurement_fields(
        self, markdown: str, filename: str, import_id: str, llm_fn: object,
    ) -> int:
        """使用 LLM 从 PPT 内容中提取结构化采购字段"""
        prompt = EXTRACTION_PROMPT.format(content=markdown[:8000])  # 限制长度
        try:
            # llm_fn 应返回 LLM 的文本响应
            response_text = await llm_fn(prompt)  # type: ignore
            # 解析 JSON
            items = json.loads(response_text)
            if not isinstance(items, list):
                items = [items]
        except (json.JSONDecodeError, Exception) as e:
            logger.warning(f"LLM 提取采购字段失败: {e}")
            return 0

        # 写入 SQLite
        count = 0
        for item in items:
            try:
                await storage.add_procurement_record({
                    "source_file": filename,
                    "source_file_id": import_id,
                    "category": item.get("category", "未分类"),
                    "item_name": item.get("item_name", "未知"),
                    "supplier": item.get("supplier"),
                    "unit_price": item.get("unit_price"),
                    "quantity": item.get("quantity"),
                    "total_price": item.get("total_price"),
                    "currency": item.get("currency", "CNY"),
                    "procurement_date": item.get("procurement_date"),
                    "contract_terms": json.dumps(
                        {"delivery_days": item.get("delivery_days"),
                         "payment_terms": item.get("payment_terms")},
                        ensure_ascii=False,
                    ),
                    "raw_text": item.get("raw_text", ""),
                    "confidence": 0.8,
                })
                count += 1
            except Exception as e:
                logger.warning(f"写入采购记录失败: {e}")
        return count

    async def _vectorize_chunks(
        self, chunks: List[dict], import_id: str, embedding_fn: Any,
    ) -> int:
        """将 PPT 内容分块并向量化写入 LanceDB"""
        if not chunks:
            return 0

        # 批量 Embedding（每批 32 条）
        batch_size = 32
        all_records: List[dict] = []
        texts = [c["content"] for c in chunks]

        for i in range(0, len(texts), batch_size):
            batch = texts[i:i + batch_size]
            try:
                vectors = await embedding_fn(batch)
                for j, vec in enumerate(vectors):
                    chunk = chunks[i + j]
                    all_records.append({
                        "id": chunk["id"],
                        "source_file": chunk["source_file"],
                        "slide_index": chunk["slide_index"],
                        "chunk_type": chunk["chunk_type"],
                        "content": chunk["content"],
                        "vector": vec,
                        "metadata": self._build_chunk_metadata(import_id, chunk),
                        "created_at": datetime.now(timezone.utc).isoformat(),
                    })
            except Exception as e:
                logger.warning(f"Embedding 批次失败 (batch {i}): {e}")

        if not all_records:
            return 0

        # 写入 LanceDB
        try:
            table_name = "doc_chunks"
            if table_name in self._lance_db.table_names():
                table = self._lance_db.open_table(table_name)
                table.add(all_records)
            else:
                table = self._lance_db.create_table(table_name, data=all_records)
            self._chunks_table = table
            logger.info(f"写入 {len(all_records)} 条向量到 LanceDB")
        except Exception as e:
            logger.error(f"LanceDB 写入失败: {e}")
            return 0

        return len(all_records)

    def _split_into_chunks(self, ppt_data: dict, filename: str) -> List[dict]:
        """将 PPT 数据按页分块（每页文本/表格/备注各为一个 chunk）"""
        chunks: List[dict] = []
        chunk_index = 1
        for slide in ppt_data.get("slides", []):
            slide_idx = slide.get("index", 0)

            # 文本 chunk
            text_parts = slide.get("texts", [])
            if text_parts:
                combined = "\n".join(text_parts)
                if combined.strip():
                    chunks.append({
                        "id": gen_id(),
                        "source_file": filename,
                        "slide_index": slide_idx,
                        "chunk_type": "text",
                        "chunk_index": chunk_index,
                        "content": combined[:2000],  # 限制长度
                    })
                    chunk_index += 1

            # 表格 chunk
            for table in slide.get("tables", []):
                md = table.get("markdown", "")
                if md.strip():
                    chunks.append({
                        "id": gen_id(),
                        "source_file": filename,
                        "slide_index": slide_idx,
                        "chunk_type": "table",
                        "chunk_index": chunk_index,
                        "content": md[:2000],
                    })
                    chunk_index += 1

            # 备注 chunk
            notes = slide.get("notes", "")
            if notes.strip():
                chunks.append({
                    "id": gen_id(),
                    "source_file": filename,
                    "slide_index": slide_idx,
                    "chunk_type": "note",
                    "chunk_index": chunk_index,
                    "content": notes[:2000],
                })
                chunk_index += 1
        return chunks

    async def vector_search(
        self, query_vector: List[float], limit: int = 5,
    ) -> List[dict]:
        """在 LanceDB 中进行向量语义检索"""
        if not self._lance_db:
            return []
        try:
            table_name = "doc_chunks"
            if table_name not in self._lance_db.table_names():
                return []
            table = self._lance_db.open_table(table_name)
            results = table.search(query_vector).limit(limit).to_list()
            # 转换为标准 dict 列表
            normalized_results: list[dict] = []
            for record in results:
                metadata = self._parse_chunk_metadata(record.get("metadata"))
                locator = metadata.get("locator") if isinstance(metadata.get("locator"), dict) else {}
                table_meta = metadata.get("table") if isinstance(metadata.get("table"), dict) else {}
                normalized_results.append({
                    "id": record.get("id", ""),
                    "content": record.get("content", ""),
                    "source_file": record.get("source_file", ""),
                    "slide_index": record.get("slide_index", 0),
                    "chunk_type": record.get("chunk_type", ""),
                    "chunk_index": metadata.get("chunk_index"),
                    "char_start": metadata.get("char_start"),
                    "char_end": metadata.get("char_end"),
                    "page": locator.get("page"),
                    "sheet": locator.get("sheet"),
                    "row_start": locator.get("row_start"),
                    "row_end": locator.get("row_end"),
                    "story": locator.get("story"),
                    "source": locator.get("source"),
                    "ocr_segment_index": locator.get("ocr_segment_index"),
                    "table_title": table_meta.get("title"),
                    "score": record.get("_distance", 0.0),
                })
            return normalized_results
        except Exception as e:
            logger.warning(f"向量检索失败: {e}")
            return []

    def _split_text_into_chunks(self, text: str, filename: str, chunk_size: int = 500) -> List[dict]:
        """将纯文本按固定大小分块，步长 400，重叠 100。"""
        chunks: List[dict] = []
        step = chunk_size - 100
        for index, i in enumerate(range(0, len(text), step), 1):
            raw_piece = text[i:i + chunk_size]
            piece = raw_piece.strip()
            if piece:
                leading_ws = len(raw_piece) - len(raw_piece.lstrip())
                trailing_ws = len(raw_piece) - len(raw_piece.rstrip())
                char_start = i + leading_ws
                char_end = i + len(raw_piece) - trailing_ws
                chunks.append({
                    "id": gen_id(),
                    "source_file": filename,
                    "slide_index": 0,
                    "chunk_type": "text",
                    "chunk_index": index,
                    "char_start": char_start,
                    "char_end": char_end,
                    "content": piece,
                })
        return chunks

    async def _vectorize_text_chunks(
        self, chunks: List[dict], import_id: str, embedding_fn: Any,
    ) -> int:
        """将文本块向量化写入 LanceDB。"""
        return await self._vectorize_chunks(chunks, import_id, embedding_fn)

    async def _ingest_generic(
        self, file_content: bytes, filename: str, text: str, file_type: str,
        embedding_fn: Optional[Any] = None,
        owner_id: Optional[str] = None,
    ) -> dict:
        """通用文件导入：记录到 ppt_imports，并把文本块写入 SQLite / LanceDB。"""
        import hashlib

        file_hash = hashlib.md5(file_content).hexdigest()
        import_id = await storage.record_ppt_import(
            file_name=filename,
            file_hash=file_hash,
            file_size=len(file_content),
            slide_count=0,
            owner_id=owner_id,
        )
        existing = await storage._fetchone(
            "SELECT import_status FROM ppt_imports WHERE id=?",
            (import_id,),
        )
        existing_chunk_count = await storage.count_knowledge_chunks(import_id=import_id)
        if existing and existing["import_status"] == "completed" and existing_chunk_count > 0:
            return {
                "import_id": import_id,
                "status": "duplicate",
                "file_type": file_type,
                "extracted_count": 0,
                "chunks_count": existing_chunk_count,
                "char_count": len(text),
            }

        await storage.update_ppt_import_status(import_id, "processing")

        chunks = self._split_text_into_chunks(text, filename) if text else []
        await self._reset_import_index(import_id, filename)

        stored_chunks_count = await self._persist_chunks(
            chunks, import_id, filename, file_type,
        )
        vector_chunks_count = 0
        if embedding_fn and text and self._lance_db:
            vector_chunks_count = await self._vectorize_text_chunks(chunks, import_id, embedding_fn)

        await storage.update_ppt_import_status(import_id, "completed", 0)

        chunks_count = vector_chunks_count if vector_chunks_count else stored_chunks_count
        char_count = len(text)
        logger.info(f"{file_type} 文件已导入: {filename} ({char_count} 字符, {chunks_count} 个块)")
        return {
            "import_id": import_id,
            "status": "completed",
            "file_type": file_type,
            "extracted_count": 0,
            "stored_chunks_count": stored_chunks_count,
            "vector_chunks_count": vector_chunks_count,
            "chunks_count": chunks_count,
            "char_count": char_count,
        }

    async def _ingest_pdf(self, file_content: bytes, filename: str, embedding_fn: Optional[Any] = None, owner_id: Optional[str] = None) -> dict:
        """导入 PDF 文件（同步解析放入线程池）"""
        import asyncio

        text = await asyncio.to_thread(self._extract_pdf_text_sync, file_content, filename)
        return await self._ingest_generic(file_content, filename, text, "pdf", embedding_fn, owner_id=owner_id)

    async def _ingest_docx(self, file_content: bytes, filename: str, embedding_fn: Optional[Any] = None, owner_id: Optional[str] = None) -> dict:
        """导入 DOCX 文件（同步解析放入线程池）"""
        import asyncio

        def _parse() -> str:
            try:
                import docx
                import io
                doc = docx.Document(io.BytesIO(file_content))
                return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
            except ImportError:
                logger.warning("python-docx 未安装，DOCX 文本提取不可用。仅记录文件。")
                return f"[DOCX 文件: {filename}，需安装 python-docx 以提取文本]"
            except Exception as e:
                logger.warning(f"DOCX 解析失败: {e}")
                return f"[DOCX 文件: {filename}，解析失败: {e}]"

        text = await asyncio.to_thread(_parse)
        return await self._ingest_generic(file_content, filename, text, "docx", embedding_fn, owner_id=owner_id)

    async def _ingest_image(self, file_content: bytes, filename: str, owner_id: Optional[str] = None) -> dict:
        """导入图片文件（记录文件信息，OCR 待后续集成）"""
        text = f"[图片文件: {filename}，大小: {len(file_content) / 1024:.1f}KB]"
        return await self._ingest_generic(file_content, filename, text, "image", owner_id=owner_id)

    async def _ingest_text(self, file_content: bytes, filename: str, embedding_fn: Optional[Any] = None, owner_id: Optional[str] = None) -> dict:
        """导入纯文本文件"""
        text = ""
        for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
            try:
                text = file_content.decode(encoding)
                break
            except (UnicodeDecodeError, LookupError):
                continue
        if not text:
            text = file_content.decode("utf-8", errors="replace")
        return await self._ingest_generic(file_content, filename, text, "text", embedding_fn, owner_id=owner_id)

    async def _ingest_excel(self, file_content: bytes, filename: str, embedding_fn: Optional[Any] = None, owner_id: Optional[str] = None) -> dict:
        """导入 Excel 文件（同步解析放入线程池）"""
        import asyncio

        def _parse() -> str:
            try:
                import openpyxl
                import io
                wb = openpyxl.load_workbook(io.BytesIO(file_content), read_only=True)
                parts = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    parts.append(f"## Sheet: {sheet}")
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        parts.append(" | ".join(cells))
                wb.close()
                return "\n".join(parts)
            except ImportError:
                logger.warning("openpyxl 未安装，Excel 文本提取不可用。仅记录文件。")
                return f"[Excel 文件: {filename}，需安装 openpyxl 以提取文本]"
            except Exception as e:
                logger.warning(f"Excel 解析失败: {e}")
                return f"[Excel 文件: {filename}，解析失败: {e}]"

        text = await asyncio.to_thread(_parse)
        return await self._ingest_generic(file_content, filename, text, "excel", embedding_fn, owner_id=owner_id)

    async def extract_text_fast(self, file_content: bytes, filename: str) -> dict:
        """
        只提取文件文本内容，不写入知识库。
        用于 📎 附件模式：将文件内容作为上下文发送给 LLM。
        同步文件解析操作通过 asyncio.to_thread 放入线程池，避免阻塞事件循环。
        """
        import asyncio
        ext = os.path.splitext(filename)[1].lower()

        if ext in self.PPT_EXTS:
            text = await asyncio.to_thread(self._extract_ppt_text_fast_sync, file_content, filename)
        elif ext in (".pdf",):
            text = await asyncio.to_thread(self._extract_pdf_text_sync, file_content, filename)
        elif ext in (".doc", ".docx"):
            def _parse_docx() -> str:
                try:
                    import docx as docx_module
                    import io
                    doc = docx_module.Document(io.BytesIO(file_content))
                    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
                except Exception as e:
                    return f"[DOCX 解析失败: {e}]"
            text = await asyncio.to_thread(_parse_docx)
        elif ext in (".xls", ".xlsx"):
            def _parse_excel() -> str:
                try:
                    import openpyxl
                    import io
                    wb = openpyxl.load_workbook(io.BytesIO(file_content), read_only=True)
                    parts = []
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        parts.append(f"## Sheet: {sheet}")
                        for row in ws.iter_rows(values_only=True):
                            cells = [str(c) if c is not None else "" for c in row]
                            parts.append(" | ".join(cells))
                    wb.close()
                    return "\n".join(parts)
                except Exception as e:
                    return f"[Excel 解析失败: {e}]"
            text = await asyncio.to_thread(self._extract_excel_text_fast_sync, file_content, filename)
        elif ext in self.IMAGE_EXTS:
            text = f"[图片文件: {filename}，大小: {len(file_content) / 1024:.1f}KB，暂不支持文本提取]"
        else:
            # 纯文本（轻量操作，无需线程池）
            text = ""
            for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
                try:
                    text = file_content.decode(encoding)
                    break
                except (UnicodeDecodeError, LookupError):
                    continue
            if not text:
                text = file_content.decode("utf-8", errors="replace")

        return {
            "filename": filename,
            "file_type": ext.lstrip("."),
            "text": text,
            "char_count": len(text),
        }

    async def _parse_document(self, file_content: bytes, filename: str) -> ParsedDocument:
        return await document_parser_registry.parse(file_content, filename)

    @staticmethod
    def _canonical_file_type(file_type: str) -> str:
        normalized = file_type.lower().lstrip(".")
        if normalized in {"ppt", "pptx"}:
            return "ppt"
        if normalized in {"doc", "docx"}:
            return "docx"
        if normalized in {"xls", "xlsx", "csv"}:
            return "excel"
        if normalized in {"txt", "md", "json", "xml"}:
            return "text"
        if normalized in {"png", "jpg", "jpeg", "gif", "bmp", "webp"}:
            return "image"
        if normalized == "pdf":
            return "pdf"
        return normalized or "text"

    def _collect_document_stats(self, parsed_document: ParsedDocument) -> dict:
        full_text = render_document_for_prompt(parsed_document)
        slides = {block.slide for block in parsed_document.blocks if block.slide}
        images = [block for block in parsed_document.blocks if block.block_type == "image"]
        return {
            "file_type": self._canonical_file_type(parsed_document.file_type),
            "slide_count": int(parsed_document.metadata.get("slide_count") or len(slides) or 0),
            "text_length": len(full_text),
            "char_count": len(full_text),
            "table_count": len(parsed_document.tables),
            "image_count": len(images),
        }

    async def _ingest_parsed_document(
        self,
        *,
        file_content: bytes,
        filename: str,
        parsed_document: ParsedDocument,
        llm_fn: Optional[object],
        embedding_fn: Optional[object],
        owner_id: Optional[str],
    ) -> dict:
        import hashlib

        stats = self._collect_document_stats(parsed_document)
        full_text = render_document_for_prompt(parsed_document)
        chunks = chunk_parsed_document(parsed_document)
        for chunk in chunks:
            chunk["source_file"] = filename

        file_hash = hashlib.md5(file_content).hexdigest()
        import_id = await storage.record_ppt_import(
            file_name=filename,
            file_hash=file_hash,
            file_size=len(file_content),
            slide_count=stats["slide_count"],
            owner_id=owner_id,
        )
        existing = await storage._fetchone(
            "SELECT import_status, slide_count FROM ppt_imports WHERE id=?",
            (import_id,),
        )
        existing_chunk_count = await storage.count_knowledge_chunks(import_id=import_id)
        if existing and existing["import_status"] == "completed" and existing_chunk_count > 0:
            return {
                "import_id": import_id,
                "status": "duplicate",
                **stats,
                "warnings": parsed_document.warnings,
                "extracted_count": 0,
                "chunks_count": existing_chunk_count,
            }

        await storage.update_ppt_import_status(import_id, "processing")
        await self._reset_import_index(
            import_id,
            filename,
            clear_procurement=stats["file_type"] == "ppt",
        )
        await storage.db.execute(
            "UPDATE ppt_imports SET slide_count=? WHERE id=?",
            (stats["slide_count"], import_id),
        )
        await storage.db.commit()

        extracted_count = 0
        if llm_fn and stats["file_type"] == "ppt" and full_text:
            extracted_count = await self._extract_procurement_fields(
                full_text,
                filename,
                import_id,
                llm_fn,
            )

        stored_chunks_count = await self._persist_chunks(
            chunks,
            import_id,
            filename,
            stats["file_type"],
        )
        vector_chunks_count = 0
        if embedding_fn and self._lance_db:
            vector_chunks_count = await self._vectorize_chunks(chunks, import_id, embedding_fn)

        await storage.update_ppt_import_status(import_id, "completed", extracted_count)
        return {
            "import_id": import_id,
            "status": "completed",
            **stats,
            "warnings": parsed_document.warnings,
            "extracted_count": extracted_count,
            "stored_chunks_count": stored_chunks_count,
            "vector_chunks_count": vector_chunks_count,
            "chunks_count": vector_chunks_count if vector_chunks_count else stored_chunks_count,
        }

    async def ingest_file(
        self,
        file_content: bytes,
        filename: str,
        llm_fn: Optional[object] = None,
        embedding_fn: Optional[object] = None,
        owner_id: Optional[str] = None,
        is_encrypted: bool = False,
    ) -> dict:
        if is_encrypted:
            file_content = decrypt_esafenet_file(file_content, filename)
        parsed_document = await self._parse_document(file_content, filename)
        return await self._ingest_parsed_document(
            file_content=file_content,
            filename=filename,
            parsed_document=parsed_document,
            llm_fn=llm_fn,
            embedding_fn=embedding_fn,
            owner_id=owner_id,
        )

    async def ingest_ppt(
        self,
        file_content: bytes,
        filename: str,
        llm_fn: Optional[object] = None,
        embedding_fn: Optional[object] = None,
        owner_id: Optional[str] = None,
    ) -> dict:
        parsed_document = await self._parse_document(file_content, filename)
        return await self._ingest_parsed_document(
            file_content=file_content,
            filename=filename,
            parsed_document=parsed_document,
            llm_fn=llm_fn,
            embedding_fn=embedding_fn,
            owner_id=owner_id,
        )

    async def extract_text_structured(self, file_content: bytes, filename: str) -> dict:
        parsed_document = await self._parse_document(file_content, filename)
        text = render_document_for_prompt(parsed_document)
        ext = os.path.splitext(filename)[1].lower()
        return {
            "filename": filename,
            "file_type": ext.lstrip("."),
            "text": text,
            "char_count": len(text),
            "warnings": parsed_document.warnings,
        }

    async def list_imports(
        self,
        user_id: Optional[str] = None,
        group_id: Optional[str] = None,
        is_admin: bool = False,
    ) -> list:
        """列出已导入的文件，按 RBAC 可见性过滤"""
        return await storage.list_ppt_imports(user_id=user_id, group_id=group_id, is_admin=is_admin)

    async def delete_import(self, import_id: str) -> dict:
        """删除指定 import_id 的知识库记录"""
        # 先查询文件名
        row = await storage._fetchone(
            "SELECT file_name FROM ppt_imports WHERE id=?", (import_id,)
        )
        if not row:
            return {"deleted": False, "message": "记录不存在"}

        filename = row["file_name"]

        # 删除采购记录
        await storage.db.execute(
            "DELETE FROM procurement_records WHERE source_file_id=?", (import_id,)
        )
        # 删除导入记录
        await storage.db.execute(
            "DELETE FROM ppt_imports WHERE id=?", (import_id,)
        )
        await storage.db.commit()

        deleted_vectors = await self._delete_vectors_for_file(filename)

        logger.info(f"已删除知识库记录: {filename} (id={import_id})")
        return {"deleted": True, "filename": filename, "deleted_vectors": deleted_vectors}

    async def get_stats(self) -> dict:
        """获取知识库统计信息"""
        # SQLite 统计
        ppt_count = await storage._fetchone("SELECT COUNT(*) as cnt FROM ppt_imports")
        record_count = await storage._fetchone("SELECT COUNT(*) as cnt FROM procurement_records")
        text_chunk_count = await storage.count_knowledge_chunks()

        # LanceDB 统计
        vector_count = 0
        if self._lance_db and "doc_chunks" in (self._lance_db.table_names() or []):
            try:
                table = self._lance_db.open_table("doc_chunks")
                vector_count = table.count_rows()
            except Exception:
                pass

        # 字段名与前端 KnowledgeStats 类型对齐
        return {
            "total_ppt_imports": ppt_count["cnt"] if ppt_count else 0,
            "completed_imports": ppt_count["cnt"] if ppt_count else 0,
            "total_procurement_records": record_count["cnt"] if record_count else 0,
            "total_text_chunks": text_chunk_count,
            "total_vector_chunks": vector_count,
        }

    async def delete_by_file(self, filename: str) -> dict:
        """删除指定文件的所有知识库数据"""
        # 删除 SQLite 记录
        await storage.db.execute(
            "DELETE FROM procurement_records WHERE source_file=?", (filename,)
        )
        await storage.db.execute(
            "DELETE FROM ppt_imports WHERE file_name=?", (filename,)
        )
        await storage.db.commit()

        deleted_vectors = await self._delete_vectors_for_file(filename)

        return {"deleted_records": True, "deleted_vectors": deleted_vectors}


# 全局单例
knowledge_service = KnowledgeService()

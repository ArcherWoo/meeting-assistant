"""
PPT 解析服务
Phase 1: python-pptx 基础解析（文本、表格、备注、元数据）
Phase 2: TODO - 集成 Docling 双引擎方案（AI 布局分析 + 高精度表格提取）
"""
import io
import hashlib
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches


class PPTParser:
    """PPT 文件解析器 - 提取结构化内容"""

    async def parse(self, file_content: bytes, filename: str) -> dict:
        """
        解析 PPT 文件，返回结构化 JSON
        包含：元数据、每页幻灯片内容（文本/表格/备注）、全文 Markdown
        """
        file_hash = hashlib.md5(file_content).hexdigest()
        prs = Presentation(io.BytesIO(file_content))

        metadata = self._extract_metadata(prs, filename, len(file_content), file_hash)
        slides = self._extract_slides(prs)
        full_markdown = self._generate_markdown(slides)

        return {
            "metadata": metadata,
            "slides": slides,
            "full_markdown": full_markdown,
            "extraction_stats": {
                "total_tables": sum(len(s.get("tables", [])) for s in slides),
                "total_images": sum(len(s.get("images", [])) for s in slides),
                "parser": "python-pptx",  # Phase 2 将变为 "docling+python-pptx"
            },
        }

    def _extract_metadata(self, prs: Presentation, filename: str, size: int, file_hash: str) -> dict:
        """提取 PPT 元数据"""
        props = prs.core_properties
        return {
            "title": props.title or filename,
            "author": props.author or "",
            "created": str(props.created) if props.created else "",
            "modified": str(props.modified) if props.modified else "",
            "slide_count": len(prs.slides),
            "file_size": f"{size / 1024 / 1024:.1f}MB",
            "file_hash": file_hash,
            "parser": "python-pptx",
        }

    def _extract_slides(self, prs: Presentation) -> list[dict]:
        """逐页提取幻灯片内容"""
        slides = []
        for idx, slide in enumerate(prs.slides):
            slide_data: dict = {
                "index": idx + 1,
                "title": self._get_slide_title(slide),
                "texts": [],
                "tables": [],
                "images": [],
                "notes": "",
            }

            # 提取文本和表格
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_data["texts"].append(text)

                if shape.has_table:
                    slide_data["tables"].append(self._extract_table(shape.table))

                # 标记图片（Phase 1 仅记录存在，Phase 2 用 Docling OCR）
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    slide_data["images"].append({"desc": shape.name, "index": len(slide_data["images"])})

            # 提取备注
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                slide_data["notes"] = notes_text

            slides.append(slide_data)
        return slides

    def _get_slide_title(self, slide) -> str:
        """获取幻灯片标题"""
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        return ""

    def _extract_table(self, table) -> dict:
        """提取表格数据为行列结构"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        # 生成 Markdown 表格
        if rows:
            header = "| " + " | ".join(rows[0]) + " |"
            separator = "| " + " | ".join(["---"] * len(rows[0])) + " |"
            body = "\n".join("| " + " | ".join(row) + " |" for row in rows[1:])
            markdown = f"{header}\n{separator}\n{body}"
        else:
            markdown = ""

        return {"rows": rows, "markdown": markdown, "row_count": len(rows)}

    def _generate_markdown(self, slides: list[dict]) -> str:
        """将所有幻灯片内容合并为完整 Markdown"""
        parts = []
        for slide in slides:
            title = slide["title"] or f"幻灯片 {slide['index']}"
            parts.append(f"## {title}\n")

            for text in slide["texts"]:
                if text != slide["title"]:  # 避免标题重复
                    parts.append(text)

            for table in slide["tables"]:
                parts.append(table["markdown"])

            if slide["notes"]:
                parts.append(f"\n> 备注: {slide['notes']}")

            parts.append("")  # 空行分隔

        return "\n\n".join(parts)

